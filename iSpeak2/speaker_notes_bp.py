import os
from flask import Blueprint, request, render_template, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

# Create a Blueprint for the new feature
speaker_notes_bp = Blueprint('speaker_notes_bp', __name__, template_folder='templates')

# Set upload folder and allowed file extensions
BASE_DIR = os.path.dirname(os.path.abspath(__file__))  # Ensure paths are within the app folder
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
RESULTS_FOLDER = os.path.join(BASE_DIR, "results")
ALLOWED_EXTENSIONS = {'pptx', 'ppt'}

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(RESULTS_FOLDER, exist_ok=True)

def allowed_file(filename):
   """Check if the uploaded file has a valid extension."""
   return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_notes_to_word(pptx_path, word_output_path):
   """Extract speaker notes from a PowerPoint presentation into a Word document."""
   presentation = Presentation(pptx_path)
   document = Document()

   slide_number = 1
   note_found = False

   for slide in presentation.slides:
       if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
           note_text = slide.notes_slide.notes_text_frame.text.strip()
           note_found = True
           document.add_heading(f"Slide {slide_number}", level=2)
           paragraph = document.add_paragraph(note_text)
           paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT

           # Customize font
           for run in paragraph.runs:
               run.font.size = Pt(12)
       else:
           # If no notes, still mention the slide
           document.add_heading(f"Slide {slide_number}", level=2)
           document.add_paragraph("No notes available.")

       slide_number += 1

   # Save the Word document
   document.save(word_output_path)
   return note_found

@speaker_notes_bp.route("/", methods=["GET", "POST"])
def upload_file():
   """Render the upload page and process file uploads."""
   if request.method == "POST":
       # Check if the file is present in the request
       if "file" not in request.files:
           return "No file uploaded", 400

       file = request.files["file"]

       # Validate the file
       if file.filename == "":
           return "No file selected", 400
       if not allowed_file(file.filename):
           return "Invalid file type. Please upload a .ppt or .pptx file.", 400

       # Save the uploaded file securely
       filename = secure_filename(file.filename)
       uploaded_path = os.path.join(UPLOAD_FOLDER, filename)
       file.save(uploaded_path)

       # Generate the output Word document
       output_filename = f"{os.path.splitext(filename)[0]}_speaker_notes.docx"
       output_path = os.path.join(RESULTS_FOLDER, output_filename)

       # Check if the file already exists and replace it
       if os.path.exists(output_path):
           os.remove(output_path)

       note_found = extract_notes_to_word(uploaded_path, output_path)
       if note_found:
           return send_file(output_path, as_attachment=True, download_name=output_filename)
       else:
           return "No speaker notes found in the presentation."

   return render_template("transcript_extractor.html")